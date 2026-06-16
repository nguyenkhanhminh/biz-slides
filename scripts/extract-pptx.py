#!/usr/bin/env python3
"""
extract-pptx.py — Extract content from PowerPoint files for biz-slides conversion.

Usage:
    python scripts/extract-pptx.py <input.pptx> <output_dir>

Output:
    <output_dir>/
        content.md        — All slide text, titles, bullets, speaker notes
        assets/           — Extracted images (slide-N-img-M.png)
        manifest.json     — Machine-readable slide manifest for the agent
"""

import sys
import os
import json
import re
from pathlib import Path


def extract(pptx_path: str, output_dir: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Missing dependency. Run: pip install python-pptx")
        sys.exit(1)

    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    assets_dir = output_dir / "assets"
    output_dir.mkdir(parents=True, exist_ok=True)
    assets_dir.mkdir(exist_ok=True)

    prs = Presentation(pptx_path)
    slides_data = []
    content_lines = [f"# {pptx_path.stem}\n", f"Extracted: {pptx_path.name}\n"]

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide": slide_num,
            "title": None,
            "body": [],
            "notes": None,
            "images": [],
        }

        content_lines.append(f"\n---\n\n## Slide {slide_num}")

        # Extract text from all text frames
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Heuristic: first non-empty text on slide is likely the title
                    if slide_info["title"] is None and len(text) < 120:
                        slide_info["title"] = text
                        content_lines.append(f"\n**{text}**")
                    else:
                        # Detect bullet level from paragraph indent level
                        indent = para.level or 0
                        prefix = "  " * indent + "- "
                        texts.append(text)
                        slide_info["body"].append({"level": indent, "text": text})
                        content_lines.append(f"\n{prefix}{text}")

        # Extract images
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    ext = image.ext
                    img_filename = f"slide-{slide_num}-img-{shape_idx}.{ext}"
                    img_path = assets_dir / img_filename
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    slide_info["images"].append(img_filename)
                    content_lines.append(f"\n[Image: {img_filename}]")
                except Exception as e:
                    content_lines.append(f"\n[Image: extraction failed — {e}]")

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_info["notes"] = notes_text
                content_lines.append(f"\n\n> Notes: {notes_text}")

        slides_data.append(slide_info)

    # Write content.md
    with open(output_dir / "content.md", "w", encoding="utf-8") as f:
        f.write("\n".join(content_lines))

    # Write manifest.json
    manifest = {
        "source": pptx_path.name,
        "slide_count": len(slides_data),
        "has_images": any(s["images"] for s in slides_data),
        "slides": slides_data,
    }
    with open(output_dir / "manifest.json", "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)

    print(f"✓ Extracted {len(slides_data)} slides")
    print(f"  Content:  {output_dir / 'content.md'}")
    print(f"  Manifest: {output_dir / 'manifest.json'}")
    if any(s["images"] for s in slides_data):
        total_images = sum(len(s["images"]) for s in slides_data)
        print(f"  Images:   {total_images} saved to {assets_dir}/")
    print("\nSlide summary:")
    for s in slides_data:
        title = s["title"] or "(no title)"
        imgs = f" [{len(s['images'])} img]" if s["images"] else ""
        notes = " [notes]" if s["notes"] else ""
        print(f"  {s['slide']:2d}. {title}{imgs}{notes}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python scripts/extract-pptx.py <input.pptx> <output_dir>")
        sys.exit(1)
    extract(sys.argv[1], sys.argv[2])
